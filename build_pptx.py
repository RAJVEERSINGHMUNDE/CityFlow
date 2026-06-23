"""
build_pptx.py
─────────────
Generates CityFlow.pptx from PRESENTATION.md.

Usage:
    python build_pptx.py
Output:
    CityFlow.pptx in the repo root.
"""
import re
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

REPO_ROOT = Path(__file__).parent
SRC_MD    = REPO_ROOT / "PRESENTATION.md"
OUT_PPTX  = REPO_ROOT / "CityFlow.pptx"

# ── Brand palette (matches the new dashboard) ────────────────────────────
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_50  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

BLUE_700  = RGBColor(0x1D, 0x4E, 0xD8)
BLUE_600  = RGBColor(0x25, 0x63, 0xEB)
BLUE_50   = RGBColor(0xEF, 0xF6, 0xFF)

ROSE_600  = RGBColor(0xE1, 0x1D, 0x48)
AMBER_600 = RGBColor(0xD9, 0x77, 0x06)
EMERALD_600 = RGBColor(0x05, 0x96, 0x69)

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout

# ── helpers ──────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=SLATE_900,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri",
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb, tf, p, run

def add_multiline(slide, x, y, w, h, lines, *, size=14, color=SLATE_900,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.25,
                  font="Calibri"):
    """lines: list of dicts {text, bold?, size?, color?, bullet?, indent?}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        p.line_spacing = ln.get("line_spacing", line_spacing)
        if ln.get("space_before"):
            p.space_before = Pt(ln["space_before"])
        if ln.get("space_after"):
            p.space_after = Pt(ln["space_after"])
        text = ln["text"]
        bullet = ln.get("bullet", False)
        if bullet:
            text = "•  " + text
        indent = ln.get("indent", 0)
        if indent:
            p.level = indent
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(ln.get("size", size))
        run.font.bold = ln.get("bold", False)
        run.font.italic = ln.get("italic", False)
        run.font.color.rgb = ln.get("color", color)
    return tb

def page_chrome(slide, page_num, total, section):
    """Header bar + page number footer."""
    # top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), BLUE_700)
    # section label top-right
    add_text(slide, Inches(8.5), Inches(0.32), Inches(4.6), Inches(0.3),
             section.upper(), size=9, bold=True, color=SLATE_500, align=PP_ALIGN.RIGHT)
    # footer divider
    add_rect(slide, Inches(0.6), Inches(7.10), Inches(12.13), Emu(9525), SLATE_200)
    # footer text
    add_text(slide, Inches(0.6), Inches(7.18), Inches(6), Inches(0.25),
             "CityFlow — Event-Driven Congestion Intelligence for Bengaluru",
             size=9, color=SLATE_500)
    add_text(slide, Inches(7.0), Inches(7.18), Inches(5.83), Inches(0.25),
             f"{page_num} / {total}", size=9, color=SLATE_500, align=PP_ALIGN.RIGHT)

def title_block(slide, kicker, title):
    """Standard slide title with a small blue kicker chip above it."""
    # kicker
    add_rect(slide, Inches(0.6), Inches(0.55), Inches(0.18), Inches(0.18), BLUE_700)
    add_text(slide, Inches(0.88), Inches(0.45), Inches(8), Inches(0.32),
             kicker.upper(), size=11, bold=True, color=BLUE_700)
    # title
    add_text(slide, Inches(0.6), Inches(0.85), Inches(12.13), Inches(0.7),
             title, size=30, bold=True, color=SLATE_900)
    # subtitle underline
    add_rect(slide, Inches(0.6), Inches(1.62), Inches(0.6), Emu(28575), BLUE_700)

# ── Slide builders ───────────────────────────────────────────────────────

def slide_title():
    s = prs.slides.add_slide(BLANK)
    # full bleed background
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, SLATE_900)
    # subtle decorative line
    add_rect(s, Inches(0.8), Inches(2.6), Inches(0.8), Inches(0.06), BLUE_700)
    # small label
    add_text(s, Inches(0.8), Inches(2.2), Inches(8), Inches(0.4),
             "HACKATHON SUBMISSION  ·  PROBLEM STATEMENT 2",
             size=11, bold=True, color=BLUE_600)
    # main title
    add_text(s, Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.2),
             "CityFlow", size=72, bold=True, color=WHITE)
    # subtitle
    add_text(s, Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.6),
             "Event-Driven Congestion Intelligence for Bengaluru",
             size=24, color=SLATE_300)
    # tagline
    add_text(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.5),
             "A Graph-AI digital twin that turns 8,205 historical traffic events",
             size=15, color=SLATE_500, italic=True)
    add_text(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(0.5),
             "into a learnable, data-driven playbook for every rally, breakdown, and festival.",
             size=15, color=SLATE_500, italic=True)
    # four pillars tag row
    pillars = ["FORECAST", "ROUTE", "DEPLOY", "LEARN"]
    x = Inches(0.8)
    for p in pillars:
        w = Inches(1.45)
        add_rect(s, x, Inches(6.6), w, Inches(0.4), SLATE_800)
        add_text(s, x, Inches(6.6), w, Inches(0.4),
                 p, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += w + Inches(0.1)
    # bottom strip
    add_rect(s, Inches(0), Inches(7.30), SLIDE_W, Inches(0.20), BLUE_700)

def slide_section_divider(num, title, subtitle):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, SLATE_900)
    add_rect(s, Inches(0.8), Inches(2.7), Inches(0.8), Inches(0.06), BLUE_700)
    add_text(s, Inches(0.8), Inches(2.3), Inches(8), Inches(0.4),
             f"PART {num}", size=12, bold=True, color=BLUE_600)
    add_text(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.4),
             title, size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.8),
             subtitle, size=18, color=SLATE_300)
    add_rect(s, Inches(0), Inches(7.30), SLIDE_W, Inches(0.20), BLUE_700)

def slide_problem():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 3, 38, "Part 1 · Foundations")
    title_block(s, "The Problem", "Event-driven congestion is unquantified, unplanned, and unforgotten")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(4.5), [
        {"text": "Political rallies. Festivals. Sports events. Construction. Sudden breakdowns. Every day in Bengaluru, the city loses hours of productivity because we cannot answer three questions:", "size": 16, "color": SLATE_700, "space_after": 12},
        {"text": "How bad will this event be?", "size": 18, "bold": True, "color": BLUE_700, "bullet": True, "space_after": 4},
        {"text": "Where should the police go?", "size": 18, "bold": True, "color": BLUE_700, "bullet": True, "space_after": 4},
        {"text": "Which routes should drivers take?", "size": 18, "bold": True, "color": BLUE_700, "bullet": True, "space_after": 18},
        {"text": "Today, all three are answered by experience and gut feel. Nobody learns from the last event.", "size": 16, "color": SLATE_700, "italic": True, "space_after": 8},
        {"text": "CityFlow turns 8,205 historical Bengaluru traffic incidents into a learnable, mathematically grounded playbook for every event that disrupts the city.", "size": 16, "color": SLATE_900, "bold": True},
    ])

def slide_data():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 4, 38, "Part 1 · Foundations")
    title_block(s, "The Data", "dataset/2.csv — the city's own incident log")

    # 6 stat tiles in 2 rows
    tiles = [
        ("8,205", "total events"),
        ("151", "days of data"),
        ("7,706", "unplanned (94.3%)"),
        ("467", "planned (5.7%)"),
        ("870", "Kannada descriptions"),
        ("1,007", "censored (right-censored)"),
    ]
    x0, y0 = Inches(0.6), Inches(2.05)
    tw, th, gap = Inches(2.0), Inches(1.5), Inches(0.16)
    for i, (big, small) in enumerate(tiles):
        col = i % 6
        x = x0 + (tw + gap) * col
        add_rect(s, x, y0, tw, th, WHITE, line=SLATE_200)
        add_rect(s, x, y0, tw, Inches(0.06), BLUE_700)
        add_text(s, x, y0 + Inches(0.25), tw, Inches(0.7),
                 big, size=28, bold=True, color=SLATE_900, align=PP_ALIGN.CENTER)
        add_text(s, x, y0 + Inches(0.95), tw, Inches(0.4),
                 small, size=11, color=SLATE_500, align=PP_ALIGN.CENTER)

    # Headline insight
    add_rect(s, Inches(0.6), Inches(3.85), Inches(12.13), Inches(1.1), BLUE_50)
    add_text(s, Inches(0.85), Inches(3.95), Inches(11.6), Inches(0.4),
             "Headline insight from the raw data", size=11, bold=True, color=BLUE_700)
    add_text(s, Inches(0.85), Inches(4.25), Inches(11.6), Inches(0.6),
             "BMTC buses take 77.5 min to clear · private cars take 38.5 min — a 2× gap, visible in the raw data without any modelling.",
             size=16, bold=True, color=SLATE_900)

    # Bottom block: cox + route_path counts
    add_multiline(s, Inches(0.6), Inches(5.15), Inches(12.13), Inches(1.8), [
        {"text": "Training cohorts we will use later:", "size": 13, "bold": True, "color": SLATE_700, "space_after": 6},
        {"text": "2,983 rows with both start and close timestamps → Cox Proportional Hazards survival model", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 4},
        {"text": "137 rows with a route_path line-string → snap closures to real OSM edges for planned events", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 4},
        {"text": "6,813 free-text operator descriptions → multilingual NLP weak-label classifier", "size": 13, "bullet": True, "color": SLATE_700},
    ])

def slide_architecture():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 5, 38, "Part 1 · Foundations")
    title_block(s, "System Architecture", "Four pillars, one JSON API, every component replaceable")

    # Pipeline blocks top-down
    boxes = [
        ("dataset / 2.csv", "8,205 Bengaluru traffic events", SLATE_700, WHITE),
        ("DataPipeline", "clean · filter · separate planned vs unplanned", SLATE_100, SLATE_900),
        ("ML models (parallel)", "Severity GBM+RF  ·  Cox PH  ·  Hotspot  ·  LaBSE NLP", SLATE_100, SLATE_900),
        ("GraphEngine", "cached Bengaluru OSMnx graph  (155K nodes · 393K edges)", SLATE_100, SLATE_900),
        ("CongestionSimulator", "BPR shockwave · Dijkstra · barricades · diversion plan", BLUE_700, WHITE),
        ("ManpowerAllocator", "linear formula · refit via np.linalg.lstsq", SLATE_100, SLATE_900),
        ("Flask API  +  SQLite", "async task queue · operational memory", SLATE_800, WHITE),
        ("React + Vite Dashboard", "5-step Story view · Plan view · dark mode", SLATE_900, WHITE),
    ]
    x = Inches(0.6)
    w = Inches(12.13)
    h = Inches(0.55)
    gap = Inches(0.05)
    y = Inches(2.05)
    for i, (label, sub, bg, fg) in enumerate(boxes):
        add_rect(s, x, y, w, h, bg, line=SLATE_200)
        add_text(s, Inches(0.85), y + Inches(0.06), Inches(3.2), h,
                 label, size=12, bold=True, color=fg, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(4.1), y + Inches(0.06), w - Inches(3.5), h,
                 sub, size=11, color=fg, italic=True, anchor=MSO_ANCHOR.MIDDLE)
        # arrow
        if i < len(boxes) - 1:
            ay = y + h + Emu(0)
            add_text(s, Inches(0.6), ay, Inches(0.3), Inches(0.18),
                     "▼", size=10, color=SLATE_500, align=PP_ALIGN.CENTER)
        y += h + gap + Inches(0.13)

def slide_bpr():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 6, 38, "Part 1 · Math Foundations")
    title_block(s, "Math #1 · BPR Travel Time", "How travel time grows as a road fills up")

    # Formula box
    add_rect(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(1.2), BLUE_50)
    add_text(s, Inches(0.6), Inches(2.1), Inches(12.13), Inches(0.4),
             "BUREAU OF PUBLIC ROADS FORMULA",
             size=11, bold=True, color=BLUE_700, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.45), Inches(12.13), Inches(0.6),
             "t  =  t₀  ·  ( 1  +  α  ·  ( V / C ) ^ β )",
             size=24, bold=True, color=SLATE_900, align=PP_ALIGN.CENTER, font="Consolas")

    add_multiline(s, Inches(0.6), Inches(3.4), Inches(12.13), Inches(3.5), [
        {"text": "Where:", "size": 14, "bold": True, "color": SLATE_700, "space_after": 4},
        {"text": "t₀   = free-flow travel time on the edge", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "V   = vehicle volume (per hour)", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "C   = capacity of the edge (per hour)", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "α   = 0.15   ·   β = 4  (standard coefficients)", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 14},
        {"text": "The β = 4 exponent is what gives the curve its hockey-stick shape — travel time stays close to free-flow until V/C approaches 0.8, then explodes.", "size": 15, "color": SLATE_700, "italic": True, "space_after": 12},
        {"text": "Why a uniform scalar is useless:", "size": 14, "bold": True, "color": ROSE_600, "space_after": 4},
        {"text": "If every edge is scaled by the same number, the shortest path between two nodes is mathematically unchanged. We need per-class, per-hour differentiation (see slide on time-of-day weighting).", "size": 14, "color": SLATE_700},
    ])

def slide_shockwave():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 7, 38, "Part 1 · Math Foundations")
    title_block(s, "Math #2 · Reverse-BFS Shockwave", "Congestion propagates upstream along the road graph, not radially")

    # Two columns: Naive vs CityFlow
    cw = Inches(5.9)
    y = Inches(2.0)
    add_rect(s, Inches(0.6), y, cw, Inches(2.7), WHITE, line=SLATE_200)
    add_rect(s, Inches(0.6), y, Inches(0.18), Inches(2.7), ROSE_600)
    add_text(s, Inches(0.95), y + Inches(0.2), cw, Inches(0.4),
             "Naïve: Euclidean circle", size=15, bold=True, color=ROSE_600)
    add_multiline(s, Inches(0.95), y + Inches(0.7), cw - Inches(0.4), Inches(2.0), [
        {"text": "Blocks every road within a radius of the event.", "size": 12, "color": SLATE_700, "space_after": 6},
        {"text": "A parallel highway 300 m away but separated by a river is wrongly shut.", "size": 12, "color": SLATE_700, "space_after": 6},
        {"text": "A dead-end alley 200 m away is treated as 12.5 acres of gridlock.", "size": 12, "color": SLATE_700},
    ])

    add_rect(s, Inches(6.83), y, cw, Inches(2.7), WHITE, line=SLATE_200)
    add_rect(s, Inches(6.83), y, Inches(0.18), Inches(2.7), EMERALD_600)
    add_text(s, Inches(7.18), y + Inches(0.2), cw, Inches(0.4),
             "CityFlow: reverse-BFS on the graph", size=15, bold=True, color=EMERALD_600)
    add_multiline(s, Inches(7.18), y + Inches(0.7), cw - Inches(0.4), Inches(2.0), [
        {"text": "Walks incoming edges from the epicenter.", "size": 12, "color": SLATE_700, "space_after": 6},
        {"text": "Respects topology — roads only blocked if connected to the event.", "size": 12, "color": SLATE_700, "space_after": 6},
        {"text": "Each edge's capacity decays with its graph distance from the epicenter.", "size": 12, "color": SLATE_700},
    ])

    # Formula
    add_rect(s, Inches(0.6), Inches(4.9), Inches(12.13), Inches(2.0), SLATE_900)
    add_text(s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(0.4),
             "CAPACITY DECAY  (per edge at graph distance d from the epicenter)",
             size=10, bold=True, color=BLUE_600)
    add_text(s, Inches(0.85), Inches(5.35), Inches(11.6), Inches(0.5),
             "if d ≤ closure_radius (50m):         capacity ← capacity × 0.05",
             size=14, color=WHITE, font="Consolas")
    add_text(s, Inches(0.85), Inches(5.85), Inches(11.6), Inches(0.5),
             "elif d ≤ spillover_radius (1km):     capacity ← capacity × (0.1 + 0.9 · d / 1000)",
             size=14, color=WHITE, font="Consolas")
    add_text(s, Inches(0.85), Inches(6.35), Inches(11.6), Inches(0.4),
             "else:                                                                  capacity unchanged",
             size=14, color=WHITE, font="Consolas")

def slide_flow_selection():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 8, 38, "Part 1 · Math Foundations")
    title_block(s, "Math #3 · Volume-Aware Flow Selection", "Up to 3 arterial OD movements, ranked by traffic they actually carry")

    add_multiline(s, Inches(0.6), Inches(2.05), Inches(12.13), Inches(2.3), [
        {"text": "For every boundary node upstream and every boundary node downstream of the event:", "size": 15, "color": SLATE_700, "space_after": 8},
        {"text": "Run Dijkstra from origin → epicenter", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 4},
        {"text": "Run Dijkstra from epicenter → destination", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 4},
        {"text": "Concatenate, score, and rank", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 8},
        {"text": "Pick the top 3, with the constraint that no two share the same origin or destination — we never get two flows that are really the same corridor.", "size": 14, "color": SLATE_700, "italic": True},
    ])

    add_rect(s, Inches(0.6), Inches(4.5), Inches(12.13), Inches(1.2), BLUE_50)
    add_text(s, Inches(0.6), Inches(4.6), Inches(12.13), Inches(0.4),
             "SCORING FUNCTION", size=11, bold=True, color=BLUE_700, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(4.95), Inches(12.13), Inches(0.6),
             "score  =  distance_km  ×  average_road_capacity_along_route",
             size=22, bold=True, color=SLATE_900, align=PP_ALIGN.CENTER, font="Consolas")

    add_text(s, Inches(0.6), Inches(5.95), Inches(12.13), Inches(0.6),
             "Longer, more arterial routes carry more vehicles → volume proxy without needing a demand matrix.",
             size=14, color=SLATE_700, italic=True, align=PP_ALIGN.CENTER)

def slide_compliance():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 9, 38, "Part 1 · Math Foundations")
    title_block(s, "Math #4 · Police Compliance", "A barricade without police is a suggestion, not a barricade")

    add_multiline(s, Inches(0.6), Inches(2.05), Inches(12.13), Inches(1.6), [
        {"text": "Bengaluru compliance with un-manned barricades ≈ 40%.  With police ≈ 95%.", "size": 15, "color": SLATE_700, "space_after": 4},
        {"text": "We model this with a single binary flag passed through the simulator.", "size": 14, "color": SLATE_500, "italic": True},
    ])

    # Two cards
    y = Inches(3.85)
    cw = Inches(5.9)
    # Without police
    add_rect(s, Inches(0.6), y, cw, Inches(2.7), WHITE, line=SLATE_200)
    add_rect(s, Inches(0.6), y, Inches(0.18), Inches(2.7), ROSE_600)
    add_text(s, Inches(0.95), y + Inches(0.2), cw, Inches(0.4),
             "police_deployed = False", size=14, bold=True, color=ROSE_600, font="Consolas")
    add_text(s, Inches(0.95), y + Inches(0.65), cw, Inches(0.4),
             "40% compliance", size=22, bold=True, color=SLATE_900)
    add_multiline(s, Inches(0.95), y + Inches(1.2), cw - Inches(0.4), Inches(1.4), [
        {"text": "Closed edges stay in the graph", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 4},
        {"text": "Their travel time is multiplied by 1 / 0.4 = 2.5×", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 4},
        {"text": "Captures the non-compliance cost without removing the edge's existence", "size": 12, "bullet": True, "color": SLATE_700},
    ])

    # With police
    add_rect(s, Inches(6.83), y, cw, Inches(2.7), WHITE, line=SLATE_200)
    add_rect(s, Inches(6.83), y, Inches(0.18), Inches(2.7), EMERALD_600)
    add_text(s, Inches(7.18), y + Inches(0.2), cw, Inches(0.4),
             "police_deployed = True", size=14, bold=True, color=EMERALD_600, font="Consolas")
    add_text(s, Inches(7.18), y + Inches(0.65), cw, Inches(0.4),
             "95% compliance", size=22, bold=True, color=SLATE_900)
    add_multiline(s, Inches(7.18), y + Inches(1.2), cw - Inches(0.4), Inches(1.4), [
        {"text": "Closed edges are removed from the graph", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 4},
        {"text": "Dijkstra routes through alternative paths", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 4},
        {"text": "API returns both scenarios — operators see how much of the win is police presence vs. geometry alone", "size": 12, "bullet": True, "color": SLATE_700},
    ])

def slide_tod():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 10, 38, "Part 1 · Math Foundations")
    title_block(s, "Math #5 · Per-Class Time-of-Day Weighting", "Make the recommended route actually differ between peak and night")

    # Formula
    add_rect(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(1.1), BLUE_50)
    add_text(s, Inches(0.6), Inches(2.1), Inches(12.13), Inches(0.4),
             "DIFFERENTIAL WEIGHTING", size=11, bold=True, color=BLUE_700, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.45), Inches(12.13), Inches(0.6),
             "multiplier = max( 0.1,  1 + ( hourly_mult − 1 ) × sensitivity[highway] )",
             size=20, bold=True, color=SLATE_900, align=PP_ALIGN.CENTER, font="Consolas")

    # Sensitivity table
    table_data = [
        ("Highway class", "Sensitivity", "At rush (hourly_mult = 1.6)"),
        ("motorway",      "1.6",  "1 + 0.6 × 1.6 = 1.96×"),
        ("trunk",         "1.5",  "1 + 0.6 × 1.5 = 1.90×"),
        ("primary",       "1.4",  "1 + 0.6 × 1.4 = 1.84×"),
        ("secondary",     "1.2",  "1 + 0.6 × 1.2 = 1.72×"),
        ("residential",   "0.7",  "1 + 0.6 × 0.7 = 1.42×"),
    ]
    tx = Inches(0.6)
    ty = Inches(3.4)
    tw_total = Inches(12.13)
    th = Inches(0.45)
    col_w = [Inches(3.5), Inches(3.0), Inches(5.63)]
    for ri, row in enumerate(table_data):
        is_header = ri == 0
        ry = ty + th * ri
        bg = SLATE_900 if is_header else (WHITE if ri % 2 == 1 else SLATE_50)
        add_rect(s, tx, ry, tw_total, th, bg, line=SLATE_200 if not is_header else None)
        cx = tx
        for ci, cell in enumerate(row):
            color = WHITE if is_header else (SLATE_500 if ci == 0 else SLATE_900)
            bold  = is_header or ci == 1
            add_text(s, cx + Inches(0.15), ry, col_w[ci], th,
                     cell, size=12, bold=bold, color=color, anchor=MSO_ANCHOR.MIDDLE)
            cx += col_w[ci]

    add_text(s, Inches(0.6), Inches(6.5), Inches(12.13), Inches(0.5),
             "Arterials amplify the rush signal · residential roads dampen it.  At 3 AM the relationship inverts.",
             size=14, color=SLATE_700, italic=True, align=PP_ALIGN.CENTER)

def slide_barricade():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 11, 38, "Part 1 · Math Foundations")
    title_block(s, "Math #6 · Continuous-Flow Barricades", "Never trap a vehicle into a U-turn")

    # The naive flaw
    add_rect(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(1.3), ROSE_600)
    add_text(s, Inches(0.6), Inches(2.05), Inches(12.13), Inches(0.4),
             "THE NAÏVE FLAW", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.45), Inches(12.13), Inches(0.8),
             "Place a barricade at the immediate upstream node u.  If u is a dead-end, cars hit the barricade, can't turn around, and gridlock.",
             size=15, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # The fix
    add_rect(s, Inches(0.6), Inches(3.55), Inches(12.13), Inches(1.1), EMERALD_600)
    add_text(s, Inches(0.6), Inches(3.6), Inches(12.13), Inches(0.4),
             "THE FIX", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(3.95), Inches(12.13), Inches(0.6),
             "Walk upstream from each closed edge until a node has a safe exit.",
             size=15, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

    # Formula
    add_rect(s, Inches(0.6), Inches(4.9), Inches(12.13), Inches(1.0), SLATE_900)
    add_text(s, Inches(0.6), Inches(5.0), Inches(12.13), Inches(0.4),
             "SAFE OUT-DEGREE  (criteria for placing a barricade at node n)",
             size=10, bold=True, color=BLUE_600, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(5.35), Inches(12.13), Inches(0.5),
             "safe_out_degree(n)  =  | { v : (n, v) is open  and  v ≠ epicenter } |   ≥  1",
             size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Consolas")

    add_text(s, Inches(0.6), Inches(6.05), Inches(12.13), Inches(0.7),
             "If no upstream node has a safe exit, the barricade is dropped.  Every recommended barricade is then validated: it must block ≥ 1 closure entry AND offer ≥ 1 alternate exit.",
             size=13, color=SLATE_700, italic=True, align=PP_ALIGN.CENTER)

def slide_severity():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 12, 38, "Part 1 · ML Models")
    title_block(s, "ML #1 · Severity Prediction", "Two heads: continuous resolution time + discrete response level")

    # Two columns
    y = Inches(2.05)
    cw = Inches(5.9)
    add_rect(s, Inches(0.6), y, cw, Inches(4.6), WHITE, line=SLATE_200)
    add_rect(s, Inches(0.6), y, Inches(0.18), Inches(4.6), BLUE_700)
    add_text(s, Inches(0.95), y + Inches(0.2), cw, Inches(0.4),
             "Head 1 · Gradient Boosting Regressor", size=15, bold=True, color=BLUE_700)
    add_multiline(s, Inches(0.95), y + Inches(0.7), cw - Inches(0.4), Inches(3.7), [
        {"text": "Predicts log(1 + resolution_minutes) — log-space to handle right-skewed clearance times.", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "Features:", "size": 12, "bold": True, "color": SLATE_700, "space_after": 4},
        {"text": "cause, event_type, requires_closure, priority", "size": 11, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "hour_sin / hour_cos (cyclical time encoding)", "size": 11, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "day_of_week, zone, spatial_cluster (KMeans)", "size": 11, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "junction_hotspot_score (historical frequency)", "size": 11, "bullet": True, "color": SLATE_700, "space_after": 10},
        {"text": "5-fold cross-validation, R² reported for transparency.", "size": 12, "color": SLATE_700, "italic": True},
    ])

    add_rect(s, Inches(6.83), y, cw, Inches(4.6), WHITE, line=SLATE_200)
    add_rect(s, Inches(6.83), y, Inches(0.18), Inches(4.6), BLUE_700)
    add_text(s, Inches(7.18), y + Inches(0.2), cw, Inches(0.4),
             "Head 2 · Random Forest Classifier", size=15, bold=True, color=BLUE_700)
    add_multiline(s, Inches(7.18), y + Inches(0.7), cw - Inches(0.4), Inches(3.7), [
        {"text": "Three response levels, thresholded on resolution time:", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "Green:  ≤ 60 minutes", "size": 12, "bullet": True, "color": EMERALD_600, "bold": True, "space_after": 3},
        {"text": "Amber:  ≤ 480 minutes", "size": 12, "bullet": True, "color": AMBER_600, "bold": True, "space_after": 3},
        {"text": "Red:    otherwise", "size": 12, "bullet": True, "color": ROSE_600, "bold": True, "space_after": 10},
        {"text": "Scenario modifiers for planned events:", "size": 12, "bold": True, "color": SLATE_700, "space_after": 4},
        {"text": "35,000 attendance → +1.0 score bonus, ×1.35 minutes", "size": 11, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "5,000 attendance  → +0.5 score bonus, ×1.15 minutes", "size": 11, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Full closure      → +0.5 score bonus, ×1.20 minutes", "size": 11, "bullet": True, "color": SLATE_700},
    ])

def slide_cox():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 13, 38, "Part 1 · ML Models")
    title_block(s, "ML #2 · Cox Proportional Hazards", "Right-censored survival model — honours un-resolved incidents")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(2.3), [
        {"text": "Naïve: drop the 1,007 still-active rows.  Loses a third of the data and biases the mean downward.", "size": 14, "color": ROSE_600, "space_after": 6},
        {"text": "CityFlow: Cox PH with right-censoring via the lifelines library.", "size": 14, "color": SLATE_700, "space_after": 12},
        {"text": "For each row:", "size": 14, "bold": True, "color": SLATE_700, "space_after": 4},
        {"text": "T = (close_time − start_time) in minutes", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3, "font": "Consolas"},
        {"text": "E = 1 if close_time observed,  0 if censored at the observation boundary", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 8, "font": "Consolas"},
    ])

    # covariates box
    add_rect(s, Inches(0.6), Inches(4.5), Inches(7.5), Inches(2.4), SLATE_50, line=SLATE_200)
    add_text(s, Inches(0.85), Inches(4.6), Inches(7.0), Inches(0.4),
             "Covariates", size=13, bold=True, color=SLATE_700)
    add_multiline(s, Inches(0.85), Inches(4.95), Inches(7.0), Inches(2.0), [
        {"text": "vehicle type (BMTC bus, heavy, LCV)", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "corridor class (highway vs non-corridor)", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "requires_road_closure", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "hour_sin / hour_cos (cyclical)", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "L2 penalty = 0.1 for stability", "size": 12, "bullet": True, "color": SLATE_700},
    ])

    # C-index card
    add_rect(s, Inches(8.43), Inches(4.5), Inches(4.3), Inches(2.4), BLUE_50)
    add_text(s, Inches(8.43), Inches(4.6), Inches(4.3), Inches(0.4),
             "C-INDEX", size=11, bold=True, color=BLUE_700, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.43), Inches(4.95), Inches(4.3), Inches(0.9),
             "0.55 – 0.70", size=36, bold=True, color=BLUE_700, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.43), Inches(5.95), Inches(4.3), Inches(0.9),
             "Probability that the model correctly orders two random incidents by clearance time. Better than chance, not a replacement for local knowledge.",
             size=11, color=SLATE_700, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def slide_nlp():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 14, 38, "Part 1 · ML Models")
    title_block(s, "ML #3 · LaBSE Multilingual NLP", "Kannada, English, mixed — one model, one embedding space")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(2.0), [
        {"text": "LaBSE = Language-agnostic BERT Sentence Embedding.  90 MB, 768-dim, multilingual, CPU-safe.", "size": 14, "color": SLATE_700, "space_after": 6},
        {"text": "Kannada 'ನಿಧಾನ' (slow) and English 'heavy traffic' land close together in embedding space.", "size": 13, "color": SLATE_500, "italic": True, "space_after": 6},
        {"text": "We attach a logistic regression head trained on weak labels extracted from the descriptions themselves.", "size": 14, "color": SLATE_700},
    ])

    add_rect(s, Inches(0.6), Inches(4.2), Inches(6.1), Inches(2.7), WHITE, line=SLATE_200)
    add_rect(s, Inches(0.6), Inches(4.2), Inches(0.18), Inches(2.7), ROSE_600)
    add_text(s, Inches(0.95), Inches(4.3), Inches(5.5), Inches(0.4),
             "Disrupted keywords", size=13, bold=True, color=ROSE_600)
    add_text(s, Inches(0.95), Inches(4.7), Inches(5.5), Inches(2.1),
             '"slow", "closed", "blocked", "gridlock", "heavy",\n"ನಿಧಾನ" (slow), "ಸಮಸ್ಯ" (problem),\n"ನಿಂತಿದೆ" (stopped), "ಕ್ಲೋಸ್" (close)',
             size=12, color=SLATE_700, font="Consolas")

    add_rect(s, Inches(6.83), Inches(4.2), Inches(5.9), Inches(2.7), WHITE, line=SLATE_200)
    add_rect(s, Inches(6.83), Inches(4.2), Inches(0.18), Inches(2.7), EMERALD_600)
    add_text(s, Inches(7.18), Inches(4.3), Inches(5.5), Inches(0.4),
             "Contained keywords", size=13, bold=True, color=EMERALD_600)
    add_text(s, Inches(7.18), Inches(4.7), Inches(5.5), Inches(2.1),
             '"no problem", "normal",\n"moving", "clear", "cleared"',
             size=12, color=SLATE_700, font="Consolas")

def slide_manpower():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 15, 38, "Part 1 · ML Models")
    title_block(s, "ML #4 · Manpower Linear Allocator", "A learnable 5-weight function that refits from your feedback")

    add_rect(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(2.0), BLUE_50)
    add_text(s, Inches(0.6), Inches(2.1), Inches(12.13), Inches(0.4),
             "OFFICERS PER BARRICADE", size=11, bold=True, color=BLUE_700, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.5), Inches(12.13), Inches(1.4),
             "intercept\n+ w_severity  × severity_score\n+ w_attendance × (expected_attendance / 1000)\n+ w_rush_hour  × is_rush_hour\n+ w_closure    × requires_closure",
             size=14, color=SLATE_900, align=PP_ALIGN.CENTER, font="Consolas",
             line_spacing=1.3)

    add_multiline(s, Inches(0.6), Inches(4.2), Inches(12.13), Inches(2.7), [
        {"text": "Default weights (hand-tuned)", "size": 13, "bold": True, "color": SLATE_700, "space_after": 4},
        {"text": "0.5  ·  0.35  ·  0.15  ·  1.2  ·  2.0", "size": 13, "color": SLATE_500, "font": "Consolas", "space_after": 12},
        {"text": "How the weights learn", "size": 13, "bold": True, "color": SLATE_700, "space_after": 4},
        {"text": "After every 10 feedback entries, CityFlow runs np.linalg.lstsq on the operator-supplied outcomes.", "size": 13, "color": SLATE_700, "space_after": 4},
        {"text": "New weights are persisted to manpower_weights.json and used for every subsequent event.", "size": 13, "color": SLATE_700, "space_after": 4},
        {"text": "The 20th event is calibrated on your corridor, your operator, your truth — not the city-wide mean.", "size": 13, "color": BLUE_700, "bold": True, "italic": True},
    ])

def slide_learning():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 16, 38, "Part 1 · ML Models")
    title_block(s, "ML #5 · Post-Event Learning", "No catastrophic forgetting. The base set only ever grows.")

    # Two boxes
    add_rect(s, Inches(0.6), Inches(2.0), Inches(5.9), Inches(4.6), WHITE, line=SLATE_200)
    add_rect(s, Inches(0.6), Inches(2.0), Inches(0.18), Inches(4.6), ROSE_600)
    add_text(s, Inches(0.95), Inches(2.1), Inches(5.4), Inches(0.4),
             "The naïve way", size=15, bold=True, color=ROSE_600)
    add_multiline(s, Inches(0.95), Inches(2.6), Inches(5.4), Inches(3.8), [
        {"text": "Refit the logistic head on only the latest 10 feedback rows.", "size": 13, "color": SLATE_700, "space_after": 8},
        {"text": "After 50 feedback entries the model has seen the 50 most-recent operator labels and nothing else.", "size": 13, "color": SLATE_700, "space_after": 8},
        {"text": "Classic catastrophic forgetting — the original weak-labelled descriptions are silently dropped.", "size": 13, "color": ROSE_600, "italic": True},
    ])

    add_rect(s, Inches(6.83), Inches(2.0), Inches(5.9), Inches(4.6), WHITE, line=SLATE_200)
    add_rect(s, Inches(6.83), Inches(2.0), Inches(0.18), Inches(4.6), EMERALD_600)
    add_text(s, Inches(7.18), Inches(2.1), Inches(5.4), Inches(0.4),
             "CityFlow's way", size=15, bold=True, color=EMERALD_600)
    add_multiline(s, Inches(7.18), Inches(2.6), Inches(5.4), Inches(3.8), [
        {"text": "Concatenate, never replace.", "size": 13, "bold": True, "color": SLATE_900, "space_after": 6},
        {"text": "X_train = vstack(X_original_weak_labels, X_feedback)", "size": 12, "color": SLATE_700, "font": "Consolas", "space_after": 4},
        {"text": "y_train = concat(y_original_weak_labels, y_feedback)", "size": 12, "color": SLATE_700, "font": "Consolas", "space_after": 10},
        {"text": "Re-fit the classifier on the full set. Re-save the augmented matrix to nlp_model.pkl.", "size": 13, "color": SLATE_700, "space_after": 8},
        {"text": "The base set grows by exactly the number of feedback rows each cycle. No forgetting — the model gets both more accurate and more durable.", "size": 13, "color": EMERALD_600, "italic": True},
    ])

def slide_impact():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 17, 38, "Part 1 · ML Models")
    title_block(s, "Impact Forecast", "Quantify the cost of the event before it happens")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(1.6), [
        {"text": "Four outputs, before the event:", "size": 14, "bold": True, "color": SLATE_700, "space_after": 6},
        {"text": "Quantifies what the event will cost the city, not just how long the incident itself will take.", "size": 13, "color": SLATE_500, "italic": True},
    ])

    # 4 stat tiles
    tiles = [
        ("Affected vehicle count", "Sum of vehicle throughput across all identified flows, scaled by route length.", BLUE_700),
        ("Person-delay minutes", "Affected vehicles × average delay per vehicle.", BLUE_700),
        ("Queue length (metres)", "spillover_radius × √duration_hours × 0.6  →  queue grows as √duration.", BLUE_700),
        ("Area Congestion Index", "0–1 score combining closed edges, spillback, and attendance.", BLUE_700),
    ]
    y0 = Inches(3.7)
    cw = Inches(2.95)
    for i, (title, sub, color) in enumerate(tiles):
        x = Inches(0.6) + (cw + Inches(0.13)) * i
        add_rect(s, x, y0, cw, Inches(2.5), WHITE, line=SLATE_200)
        add_rect(s, x, y0, Inches(0.18), Inches(2.5), color)
        add_text(s, x + Inches(0.3), y0 + Inches(0.2), cw - Inches(0.4), Inches(0.6),
                 title, size=12, bold=True, color=SLATE_900)
        add_text(s, x + Inches(0.3), y0 + Inches(0.85), cw - Inches(0.4), Inches(1.5),
                 sub, size=11, color=SLATE_700)

    add_text(s, Inches(0.6), Inches(6.35), Inches(12.13), Inches(0.5),
             "Response tier (Green / Amber / Red) is derived from a weighted sum of ACI + total delay + attendance — not just from the cause.",
             size=12, color=SLATE_700, italic=True, align=PP_ALIGN.CENTER)

def slide_async():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 18, 38, "Part 1 · Systems")
    title_block(s, "Asynchronous Architecture", "Flask + threading + SQLite, no Celery required")

    # Sequence as boxes with arrows
    steps = [
        ("1", "Frontend", "POST /api/simulate/<id>", BLUE_700, WHITE),
        ("2", "Flask", "Generate task_id, spawn thread, return 202", SLATE_700, WHITE),
        ("3", "Background thread", "Load subgraph, BPR shockwave, evaluate, render Folium, write to tasks table", SLATE_100, SLATE_900),
        ("4", "Frontend", "Poll GET /api/status/<task_id> every 2 seconds", BLUE_700, WHITE),
        ("5", "Result", "Next poll returns { status: 'success', result_json: …, map_url: … }", EMERALD_600, WHITE),
    ]
    y = Inches(2.1)
    for num, who, what, bg, fg in steps:
        add_rect(s, Inches(0.6), y, Inches(0.6), Inches(0.85), bg)
        add_text(s, Inches(0.6), y, Inches(0.6), Inches(0.85),
                 num, size=24, bold=True, color=fg, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.3), y, Inches(2.3), Inches(0.85), SLATE_100, line=SLATE_200)
        add_text(s, Inches(1.4), y, Inches(2.1), Inches(0.85),
                 who, size=12, bold=True, color=SLATE_900, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(3.7), y, Inches(9.0), Inches(0.85), WHITE, line=SLATE_200)
        add_text(s, Inches(3.85), y, Inches(8.7), Inches(0.85),
                 what, size=13, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.95)

    add_text(s, Inches(0.6), Inches(6.9), Inches(12.13), Inches(0.4),
             "Same pattern as Celery+Redis, but single-process SQLite — perfect for a single traffic-control-room deployment, zero new infrastructure.",
             size=12, color=SLATE_700, italic=True, align=PP_ALIGN.CENTER)

def slide_demo_overview():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 19, 38, "Part 2 · UI Walkthrough")
    title_block(s, "UI Tour · What the user sees", "Every screen, in order")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(4.9), [
        {"text": "Home screen — Welcome to CityFlow, three tiles (Load demo · Create your own · Past hotspots).", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 6},
        {"text": "Sidebar — event list with badges (Planned / Unplanned / Closure / Custom), search by cause.", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 6},
        {"text": "Story view — five numbered steps: Situation → Impact → Maps → Plan → Feedback.", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 6},
        {"text": "Plan view — same data in one dense screen, no narrative cards.", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 6},
        {"text": "Help modal — three tabs: Glossary · FAQ · About, with plain-language definitions of every metric.", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 6},
        {"text": "Theme toggle — Sun / Moon icon, persisted to localStorage, no flash on reload.", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 6},
        {"text": "Map iframes — real Folium renders: original route (red), AI diversion (cyan), barricades (orange).", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 6},
        {"text": "Past hotspots — historical heatmap with top-20 junctions overlaid.", "size": 14, "bullet": True, "color": SLATE_700, "space_after": 6},
        {"text": "Feedback panel — operator logs actual outcome, retrain fires every 10 entries.", "size": 14, "bullet": True, "color": SLATE_700},
    ])

def slide_screen_situation():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 20, 38, "Part 2 · UI Walkthrough")
    title_block(s, "Story view · Step 1", "What is happening — Situation card")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(5.0), [
        {"text": "Hero section", "size": 14, "bold": True, "color": BLUE_700, "space_after": 6},
        {"text": "Blue flag icon (40px rounded square, blue-50 bg, blue-700 ring)", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Event cause in 18px bold, with Planned / Unplanned / Closure pills", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Metadata line: 📅 datetime  ·  👥 attendance  ·  📍 roads_affected", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 10},
        {"text": "Severity pill (top-right)", "size": 14, "bold": True, "color": BLUE_700, "space_after": 6},
        {"text": "Green: 'All clear'  ·  Amber: 'Heads up'  ·  Red: 'Urgent action'", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 10},
        {"text": "Headline + resolution label + model confidence note", "size": 14, "bold": True, "color": BLUE_700, "space_after": 6},
        {"text": "'This event will cause moderate disruption. Expected to take about ~2.0 hrs to clear.'", "size": 12, "color": SLATE_500, "italic": True, "space_after": 10},
        {"text": "Radio / log report with NLP-flagged words highlighted in rose", "size": 14, "bold": True, "color": BLUE_700, "space_after": 6},
        {"text": "87% likely to disrupt traffic — output of the LaBSE classifier, paired with the highlighted tokens.", "size": 12, "color": SLATE_500, "italic": True},
    ])

def slide_screen_impact():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 21, 38, "Part 2 · UI Walkthrough")
    title_block(s, "Story view · Step 2", "How bad is it? — Impact assessment")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(2.5), [
        {"text": "Card has a left-border accent whose colour matches the severity:", "size": 14, "color": SLATE_700, "space_after": 6},
        {"text": "rose-600 (Red)  ·  amber-600 (Amber)  ·  emerald-600 (Green)", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 10},
        {"text": "2x2 grid of stat tiles, every one with a hover tooltip", "size": 14, "bold": True, "color": BLUE_700, "space_after": 6},
    ])

    # 2x2 stat tiles
    tiles = [
        ("How long to clear", "~2.0 hrs", "Worst-case ~13h 21m", "clock", BLUE_700),
        ("People delayed", "60,000 min", "~4,000 vehicles", "people", BLUE_700),
        ("Backed-up queue", "848 m", "metres", "layers", BLUE_700),
        ("Congestion score", "23%", "0–1 index → percent", "bolt", BLUE_700),
    ]
    y0 = Inches(4.5)
    cw = Inches(5.9)
    ch = Inches(1.25)
    for i, (label, val, sub, _icon, color) in enumerate(tiles):
        col, row = i % 2, i // 2
        x = Inches(0.6) + (cw + Inches(0.33)) * col
        y = y0 + (ch + Inches(0.18)) * row
        add_rect(s, x, y, cw, ch, WHITE, line=SLATE_200)
        add_rect(s, x, y, Inches(0.15), ch, color)
        add_text(s, x + Inches(0.3), y + Inches(0.15), cw - Inches(0.4), Inches(0.35),
                 label, size=11, bold=True, color=SLATE_500)
        add_text(s, x + Inches(0.3), y + Inches(0.5), cw - Inches(0.4), Inches(0.5),
                 val, size=22, bold=True, color=SLATE_900)
        add_text(s, x + Inches(0.3), y + Inches(0.95), cw - Inches(0.4), Inches(0.3),
                 sub, size=10, color=SLATE_500, italic=True)

def slide_screen_maps():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 22, 38, "Part 2 · UI Walkthrough")
    title_block(s, "Story view · Step 3", "What does the city look like — Map comparison")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(2.5), [
        {"text": "Three tabs in a small pill strip", "size": 14, "bold": True, "color": BLUE_700, "space_after": 6},
        {"text": "Side-by-side (default) — two iframes, left = chaos, right = CityFlow plan", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Plan only — single iframe of the recommended plan", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Past hotspots — single iframe of the historical heatmap", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 10},
        {"text": "Map layers (Folium, dark CartoDB basemap)", "size": 14, "bold": True, "color": BLUE_700, "space_after": 6},
    ])

    legend = [
        ("Original route (red)", "will get stuck", ROSE_600),
        ("AI diversion (cyan)", "recommended alternative", BLUE_600),
        ("Barricade (orange)", "police checkpoint", AMBER_600),
        ("Event location (red)", "epicenter", ROSE_600),
    ]
    y0 = Inches(4.85)
    cw = Inches(2.95)
    for i, (label, sub, color) in enumerate(legend):
        x = Inches(0.6) + (cw + Inches(0.13)) * i
        add_rect(s, x, y0, cw, Inches(1.4), WHITE, line=SLATE_200)
        add_rect(s, x + Inches(0.3), y0 + Inches(0.4), Inches(0.3), Inches(0.6), color)
        add_text(s, x + Inches(0.75), y0 + Inches(0.3), cw - Inches(0.9), Inches(0.4),
                 label, size=12, bold=True, color=SLATE_900)
        add_text(s, x + Inches(0.75), y0 + Inches(0.7), cw - Inches(0.9), Inches(0.6),
                 sub, size=11, color=SLATE_500, italic=True)

def slide_screen_plan():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 23, 38, "Part 2 · UI Walkthrough")
    title_block(s, "Story view · Step 4", "What does CityFlow recommend — Resource plan")

    # Big number
    add_rect(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(1.5), BLUE_50)
    add_text(s, Inches(0.6), Inches(2.1), Inches(12.13), Inches(0.4),
             "WITH THIS PLAN, THE AVERAGE TRIP THROUGH THE AREA IS", size=11, bold=True, color=BLUE_700, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(2.45), Inches(12.13), Inches(0.7),
             "874.3 min faster", size=36, bold=True, color=EMERALD_600, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.6), Inches(3.05), Inches(12.13), Inches(0.4),
             "85.7% less delay per trip across 3 main traffic routes.",
             size=13, color=SLATE_700, italic=True, align=PP_ALIGN.CENTER)

    # Three resource tiles
    tiles = [
        ("Officers", "44", "11 per barricade", BLUE_700),
        ("Barricades", "4", "4 validated", BLUE_700),
        ("Shift", "8h", "352 officer-hrs", BLUE_700),
    ]
    y0 = Inches(3.85)
    cw = Inches(3.95)
    for i, (label, val, sub, color) in enumerate(tiles):
        x = Inches(0.6) + (cw + Inches(0.13)) * i
        add_rect(s, x, y0, cw, Inches(1.5), WHITE, line=SLATE_200)
        add_text(s, x + Inches(0.3), y0 + Inches(0.2), cw - Inches(0.4), Inches(0.3),
                 label, size=12, bold=True, color=SLATE_500)
        add_text(s, x + Inches(0.3), y0 + Inches(0.55), cw - Inches(0.4), Inches(0.7),
                 val, size=32, bold=True, color=SLATE_900)
        add_text(s, x + Inches(0.3), y0 + Inches(1.15), cw - Inches(0.4), Inches(0.3),
                 sub, size=11, color=SLATE_500, italic=True)

    add_multiline(s, Inches(0.6), Inches(5.65), Inches(12.13), Inches(1.4), [
        {"text": "Plus:", "size": 13, "bold": True, "color": SLATE_700, "space_after": 4},
        {"text": "Routes protected — list of up to 3 flow IDs, with time saved per flow", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Barricade locations — lat/lon in monospace, with validation reason", "size": 12, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Urgency note — amber-bordered callout with response-level guidance", "size": 12, "bullet": True, "color": SLATE_700},
    ])

def slide_screen_feedback():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 24, 38, "Part 2 · UI Walkthrough")
    title_block(s, "Story view · Step 5", "After it ends — Feedback panel")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(4.9), [
        {"text": "The most important step", "size": 16, "bold": True, "color": BLUE_700, "italic": True, "space_after": 12},
        {"text": "Six fields:", "size": 14, "bold": True, "color": SLATE_700, "space_after": 6},
        {"text": "Actual resolution minutes (numeric, predicted value pre-filled as placeholder)", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Observed severity dropdown — Green / Amber / Red", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Officers actually used (numeric)", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Barricades actually used (numeric)", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Diversion plan worked? — Yes / No pill radio cards", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 3},
        {"text": "Free-text notes (optional)", "size": 13, "bullet": True, "color": SLATE_700, "space_after": 14},
        {"text": "On submit → POST /api/feedback → SQLite feedback table", "size": 14, "color": SLATE_700, "space_after": 4},
        {"text": "Every 10 entries → background thread re-fits NLP + manpower weights", "size": 14, "bold": True, "color": BLUE_700, "italic": True},
    ])

def slide_screen_helper_features():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 25, 38, "Part 2 · UI Walkthrough")
    title_block(s, "Across the app", "Recurring UI primitives")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(4.9), [
        {"text": "SeverityBadge (3 colour levels × 2 sizes)", "size": 14, "bold": True, "color": BLUE_700, "space_after": 4},
        {"text": "Pill with coloured ring + dot — Green, Amber, Red.  Small (titles) and large (situation hero).", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "EventTypeBadge (Planned / Unplanned)", "size": 14, "bold": True, "color": BLUE_700, "space_after": 4},
        {"text": "Pill with calendar (violet) or triangle (orange) icon.  Same shape in the sidebar and the situation card.", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "Hint tooltips (ⓘ)", "size": 14, "bold": True, "color": BLUE_700, "space_after": 4},
        {"text": "Hover over any ⓘ icon → plain-language explanation.  No jargon left unexplained.", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "Severity left-border card accent", "size": 14, "bold": True, "color": BLUE_700, "space_after": 4},
        {"text": "Cards use a 4px coloured left border instead of a tinted background — calm, data-dense, command-centre look.", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "Glassmorphism sidebar", "size": 14, "bold": True, "color": BLUE_700, "space_after": 4},
        {"text": "Sidebar uses bg-white/75 + backdrop-blur so the map shows through subtly.  Light or dark, looks the same.", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "Dark mode toggle (Sun / Moon in the header)", "size": 14, "bold": True, "color": BLUE_700, "space_after": 4},
        {"text": "One click, every surface repaints.  Persisted to localStorage.  No flash on reload.", "size": 13, "color": SLATE_700},
    ])

def slide_demo_flow():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 26, 38, "Part 2 · UI Walkthrough")
    title_block(s, "90-second demo flow", "Click-by-click script for the video")

    steps = [
        ("0:00", "Land on the home screen", "Welcome to CityFlow · Load demo · Create your own · See past hotspots"),
        ("0:08", "Click Load demo", "Three scenarios injected, first auto-selects"),
        ("0:15", "Step 1 — Situation", "Severity pill · headline · NLP highlight in radio report"),
        ("0:25", "Step 2 — Impact", "60,000 person-delay min · 848m queue · 23% congestion · 702 nearby events"),
        ("0:40", "Step 3 — Maps", "Two iframes load · red chaos vs. green plan · cyan diversion, orange barricades"),
        ("0:55", "Step 4 — Plan", "874.3 min faster · 44 officers · 4 barricades · 8h shift · 3 routes"),
        ("1:10", "Step 5 — Feedback", "Type 180 min · pick Amber · type 42 officers · save · green toast"),
        ("1:25", "Toggle to Plan view", "Everything collapses to one dense screen"),
        ("1:30", "Click Sun icon", "One click, dark mode · every surface, every border repaints"),
        ("1:35", "Open Help", "Glossary · FAQ · About · close"),
        ("1:50", "Click 'All events'", "Back to the home screen"),
    ]
    y = Inches(2.0)
    th = Inches(0.43)
    for i, (t, what, how) in enumerate(steps):
        bg = WHITE if i % 2 == 0 else SLATE_50
        add_rect(s, Inches(0.6), y, Inches(12.13), th, bg, line=SLATE_200)
        add_text(s, Inches(0.75), y, Inches(0.9), th,
                 t, size=11, bold=True, color=BLUE_700, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
        add_text(s, Inches(1.7), y, Inches(3.4), th,
                 what, size=12, bold=True, color=SLATE_900, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.2), y, Inches(7.5), th,
                 how, size=11, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
        y += th + Inches(0.02)

def slide_api():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 27, 38, "Part 2 · API Surface")
    title_block(s, "The API", "Every screen is backed by a single JSON API")

    endpoints = [
        ("GET",  "/api/events",                 "List top events + operator scenarios"),
        ("GET",  "/api/severity/<id>",          "Severity, clearance forecast, NLP, nearby history, impact forecast"),
        ("POST", "/api/simulate/<id>",          "Async task · returns task_id"),
        ("GET",  "/api/status/<task_id>",       "Poll for simulation result"),
        ("POST", "/api/feedback",               "Record outcome · triggers retrain every 10 entries"),
        ("GET",  "/api/feedback/summary",       "Total outcomes · mean resolution error · diversion success rate"),
        ("GET",  "/api/hotspots",               "Top 20 junctions · temporal patterns · heatmap URL"),
        ("GET",  "/api/realtime/incidents",     "Historical-replay live feed (as_of=ISO)"),
        ("POST", "/api/scenarios",              "Create a what-if"),
        ("POST", "/api/scenarios/demo",         "Load the three demo events"),
    ]
    y = Inches(2.0)
    th = Inches(0.45)
    cw1, cw2, cw3 = Inches(0.9), Inches(3.6), Inches(7.63)
    # header
    add_rect(s, Inches(0.6), y, Inches(12.13), th, SLATE_900)
    add_text(s, Inches(0.75), y, cw1, th, "Method", size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.7), y, cw2, th, "Endpoint", size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.2), y, cw3, th, "What it returns", size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    y += th

    for i, (m, e, d) in enumerate(endpoints):
        bg = WHITE if i % 2 == 0 else SLATE_50
        add_rect(s, Inches(0.6), y, Inches(12.13), th, bg, line=SLATE_200)
        mcolor = EMERALD_600 if m == "GET" else BLUE_700
        add_text(s, Inches(0.75), y, cw1, th, m, size=11, bold=True, color=mcolor, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
        add_text(s, Inches(1.7), y, cw2, th, e, size=12, color=SLATE_900, anchor=MSO_ANCHOR.MIDDLE, font="Consolas")
        add_text(s, Inches(5.2), y, cw3, th, d, size=11, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)
        y += th + Inches(0.02)

def slide_limits():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 28, 38, "Part 3 · Closing")
    title_block(s, "Limits", "Be honest about what CityFlow is and is not")

    add_multiline(s, Inches(0.6), Inches(2.0), Inches(12.13), Inches(4.9), [
        {"text": "NLP requires ~2 GB RAM", "size": 14, "bold": True, "color": ROSE_600, "space_after": 4},
        {"text": "On 4 GB servers, deploy/patch_nlp.py disables it.  The system still works — disrupted_prob falls back to 0.5.", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "No live traffic feed in the prototype", "size": 14, "bold": True, "color": ROSE_600, "space_after": 4},
        {"text": "The RealtimeFeed interface is in place.  Swap in HERE / TomTom / a CCTV pipeline with a single class change.", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "First start is slow", "size": 14, "bold": True, "color": ROSE_600, "space_after": 4},
        {"text": "OSMnx downloads the Bengaluru graph on first run.  After that it's a 170 MB cache file.", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "Cox PH C-index: 0.55–0.70", "size": 14, "bold": True, "color": ROSE_600, "space_after": 4},
        {"text": "Better than chance, not a replacement for a human dispatcher's local knowledge.", "size": 13, "color": SLATE_700, "space_after": 10},
        {"text": "Not a microscopic simulator", "size": 14, "bold": True, "color": ROSE_600, "space_after": 4},
        {"text": "CityFlow is a graph-based intervention planner.  It will not tell you the exact trajectory of a specific car — that is SUMO/AIMSUN territory.", "size": 13, "color": SLATE_700},
    ])

def slide_pillars():
    s = prs.slides.add_slide(BLANK)
    page_chrome(s, 29, 38, "Part 3 · Closing")
    title_block(s, "The four pillars", "Forecast, route, deploy, learn — all real, all measured, all auditable")

    pillars = [
        ("FORECAST", "BPR reverse-BFS shockwave\nCox PH survival · LaBSE NLP\nImpact forecast before the event", BLUE_700, "src/simulator/simulator.py\nsrc/simulator/survival_model.py\nsrc/simulator/impact_forecast.py\nsrc/simulator/nlp_impact.py"),
        ("ROUTE", "Volume-aware flow selection\nPolice compliance model\nBPR + V/C penalty · Dijkstra on the adjusted graph", BLUE_700, "src/simulator/simulator.py\n(find_affected_flows · calculate_diversion)"),
        ("DEPLOY", "Continuous-flow barricades\nManpower linear formula\nValidated locations · refit from feedback", BLUE_700, "src/simulator/simulator.py\n(recommend_barricades · validate_barricades)\nsrc/simulator/manpower.py"),
        ("LEARN", "Feedback → SQLite\nEvery 10 entries → background thread\nNo catastrophic forgetting", BLUE_700, "src/api/app.py\n(/api/feedback + retrain hook)\nsrc/simulator/nlp_impact.py\nsrc/simulator/manpower.py"),
    ]
    y = Inches(2.0)
    cw = Inches(5.9)
    ch = Inches(2.4)
    for i, (name, what, color, where) in enumerate(pillars):
        col, row = i % 2, i // 2
        x = Inches(0.6) + (cw + Inches(0.33)) * col
        yy = y + (ch + Inches(0.2)) * row
        add_rect(s, x, yy, cw, ch, WHITE, line=SLATE_200)
        add_rect(s, x, yy, Inches(0.18), ch, color)
        add_text(s, x + Inches(0.3), yy + Inches(0.15), cw - Inches(0.4), Inches(0.5),
                 name, size=15, bold=True, color=color)
        add_text(s, x + Inches(0.3), yy + Inches(0.6), cw - Inches(0.4), Inches(1.0),
                 what, size=12, color=SLATE_700, line_spacing=1.25)
        add_text(s, x + Inches(0.3), yy + Inches(1.75), cw - Inches(0.4), Inches(0.6),
                 where, size=10, color=SLATE_500, italic=True, font="Consolas", line_spacing=1.3)

def slide_closing():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, SLATE_900)
    add_rect(s, Inches(0.8), Inches(2.5), Inches(0.8), Inches(0.06), BLUE_700)
    add_text(s, Inches(0.8), Inches(2.1), Inches(8), Inches(0.4),
             "THE BOTTOM LINE", size=12, bold=True, color=BLUE_600)
    add_text(s, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.6),
             "CityFlow", size=72, bold=True, color=WHITE)
    add_multiline(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.4), [
        {"text": "Takes 8,205 historical Bengaluru traffic events,", "size": 18, "color": SLATE_300, "line_spacing": 1.4},
        {"text": "the city's own road network, and the operator's own radio reports,", "size": 18, "color": SLATE_300, "line_spacing": 1.4},
        {"text": "and turns them into a data-driven playbook", "size": 18, "color": SLATE_300, "line_spacing": 1.4},
        {"text": "for the next rally, the next breakdown, the next festival.", "size": 18, "color": SLATE_300, "line_spacing": 1.4, "space_after": 24},
        {"text": "Mathematical routing intelligence — not just coloured lines on a map.", "size": 18, "color": BLUE_600, "bold": True, "italic": True},
    ])
    add_rect(s, Inches(0), Inches(7.30), SLIDE_W, Inches(0.20), BLUE_700)

# ── Build ────────────────────────────────────────────────────────────────

slide_title()                                                                                  # 1
slide_section_divider(1, "Math & Technical Foundations", "The math, the ML, the systems")     # 2
slide_problem()                                                                                 # 3
slide_data()                                                                                    # 4
slide_architecture()                                                                           # 5
slide_bpr()                                                                                     # 6
slide_shockwave()                                                                               # 7
slide_flow_selection()                                                                          # 8
slide_compliance()                                                                              # 9
slide_tod()                                                                                     # 10
slide_barricade()                                                                               # 11
slide_severity()                                                                                # 12
slide_cox()                                                                                     # 13
slide_nlp()                                                                                     # 14
slide_manpower()                                                                                # 15
slide_learning()                                                                                # 16
slide_impact()                                                                                  # 17
slide_async()                                                                                   # 18
slide_section_divider(2, "UI Walkthrough", "Every visible feature, in order")                 # 19
slide_demo_overview()                                                                           # 20
slide_screen_situation()                                                                        # 21
slide_screen_impact()                                                                           # 22
slide_screen_maps()                                                                             # 23
slide_screen_plan()                                                                             # 24
slide_screen_feedback()                                                                         # 25
slide_screen_helper_features()                                                                  # 26
slide_demo_flow()                                                                               # 27
slide_api()                                                                                     # 28
slide_section_divider(3, "Closing", "Limits, pillars, and the bottom line")                   # 29
slide_limits()                                                                                  # 30
slide_pillars()                                                                                 # 31
slide_closing()                                                                                 # 32

prs.save(str(OUT_PPTX))
print(f"Wrote {OUT_PPTX} with {len(prs.slides)} slides")
